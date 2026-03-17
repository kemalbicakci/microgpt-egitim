"""
microgpt_aciklamali.pptx dosyasına
"Dikkat Sonrası Pipeline" slaytını ekler (mimari slaytından hemen sonra, 9. slayt).
"""

import os
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
PPTX = os.path.join(HERE, 'microgpt_aciklamali.pptx')

W, H   = Inches(13.33), Inches(7.5)
DARK   = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT = RGBColor(0x16, 0x21, 0x3e)
GOLD   = RGBColor(0xe9, 0x4f, 0x37)
WHITE  = RGBColor(0xff, 0xff, 0xff)
LIGHT  = RGBColor(0xec, 0xf0, 0xf1)
MUTED  = RGBColor(0x95, 0xa5, 0xa6)
GREEN  = RGBColor(0x27, 0xae, 0x60)
BLUE   = RGBColor(0x29, 0x80, 0xb9)
PURPLE = RGBColor(0x8e, 0x44, 0xad)
ORANGE = RGBColor(0xe6, 0x7e, 0x22)

prs = Presentation(PPTX)

def bg(slide):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = DARK

def box(slide, l, t, w, h, color, border=None, lw=0):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if border:
        sh.line.color.rgb = border; sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh

def label(slide, text, l, t, w, h,
          size=14, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color

def divider(slide, top):
    box(slide, Inches(0.5), top, Inches(12.33), Inches(0.04), GOLD)

# ── Yeni slayt ──────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)

label(sl, "Dikkat Sonrası: Wo → Artık → MLP → LM Head → Softmax",
      Inches(0.5), Inches(0.18), Inches(12.5), Inches(0.62),
      size=28, bold=True, color=GOLD)
divider(sl, Inches(0.86))

# ── Pipeline adımları — sol sütun ───────────────────────────────────────────
adimlar = [
    (BLUE,   "1 — Wo Projeksiyonu",
     "4 kafanın çıktısı yan yana dizilir (4×4=16 boyut).\n"
     "Wo matrisi (16×16) hepsini tek bir vektöre harmanlıyor.\n"
     "Her kafanın öğrendikleri sentezleniyor."),

    (GREEN,  "2 — Artık Bağlantı  ( x = dikkat + x_önceki )",
     "Dikkat bloğuna girmeden önceki x direkt toplanıyor.\n"
     "Model hiçbir şey öğrenmemiş olsa bile bilgi kaybolmuyor.\n"
     "En kötü ihtimalle dikkat çıktısı ≈ 0 olur, x olduğu gibi geçer."),

    (ORANGE, "3 — MLP Bloğu  ( FC1 → ReLU → FC2 )",
     "FC1: 16 → 64 boyuta genişler  (×4)\n"
     "ReLU: negatif değerleri sıfırlar  [ -0.3 → 0.0 ]\n"
     "FC2: 64 → 16 boyuta daralır\n"
     "Dikkat 'nereye bakayım' sorusunu çözdü.\n"
     "MLP 'bu bilgiyle ne yapayım' sorusunu çözüyor."),

    (PURPLE, "4 — Artık Bağlantı  ( x = MLP + x_önceki )",
     "MLP çıktısı da x ile toplanıyor — aynı mantık.\n"
     "Gradyanlar buradan direkt erken katmanlara ulaşabiliyor."),

    (RGBColor(0xc0,0x39,0x2b), "5 — LM Head",
     "x (16 boyut)  →  lm_head matrisi (16×30)  →  30 ham skor (logit)\n"
     "Her karakter için ayrı bir sayı: 'a'→1.2  'n'→3.8  'z'→0.1"),

    (RGBColor(0x27,0x6e,0x45), "6 — Softmax → Örnekleme",
     "logitler → exp → normalize → olasılıklar  (toplamı 1.0)\n"
     "random.choices() bu olasılıklarla bir karakter seçiyor.\n"
     "Sıcaklık düşük → kesin seçim.  Sıcaklık yüksek → yaratıcı."),
]

for i, (renk, baslik, aciklama) in enumerate(adimlar):
    top = Inches(1.0) + i * Inches(1.03)
    box(sl, Inches(0.4), top, Inches(0.06), Inches(0.88), renk)
    box(sl, Inches(0.5), top, Inches(7.8), Inches(0.88), ACCENT,
        border=renk, lw=1.0)
    label(sl, baslik,
          Inches(0.65), top + Inches(0.05), Inches(7.5), Inches(0.32),
          size=13, bold=True, color=renk)
    label(sl, aciklama,
          Inches(0.65), top + Inches(0.35), Inches(7.5), Inches(0.55),
          size=11, color=LIGHT)

# ── Sağ taraf: pipeline özet akışı ─────────────────────────────────────────
box(sl, Inches(8.9), Inches(1.0), Inches(4.1), Inches(6.2), ACCENT)
label(sl, "Tam Akış", Inches(9.0), Inches(1.05), Inches(3.9), Inches(0.38),
      size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

akis = [
    (BLUE,                       "Dikkat çıktısı  →  Wo"),
    (GREEN,                      "  + x_önceki  (artık)"),
    (ORANGE,                     "FC1 → ReLU → FC2  (MLP)"),
    (PURPLE,                     "  + x_önceki  (artık)"),
    (RGBColor(0xc0,0x39,0x2b),   "lm_head  →  30 logit"),
    (RGBColor(0x27,0x6e,0x45),   "softmax  →  olasılıklar"),
    (WHITE,                      "random.choices  →  harf"),
]

for i, (renk, metin) in enumerate(akis):
    top = Inches(1.5) + i * Inches(0.75)
    box(sl, Inches(9.1), top, Inches(3.7), Inches(0.52), DARK,
        border=renk, lw=1.5)
    label(sl, metin, Inches(9.2), top + Inches(0.08),
          Inches(3.5), Inches(0.38), size=13, bold=True, color=renk)
    # ok
    if i < len(akis) - 1:
        label(sl, "↓", Inches(10.7), top + Inches(0.52),
              Inches(0.4), Inches(0.25), size=12, color=MUTED,
              align=PP_ALIGN.CENTER)

label(sl, "Bu döngü her yeni token için tekrar başa döner.",
      Inches(8.9), Inches(6.88), Inches(4.1), Inches(0.42),
      size=10, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

# ── Slaytı doğru konuma taşı (index 8 = 9. slayt) ──────────────────────────
xml_slides  = prs.slides._sldIdLst
slides_list = list(xml_slides)
new_el      = slides_list[-1]
xml_slides.remove(new_el)
xml_slides.insert(8, new_el)   # mimari diyagramından (8) hemen sonra

prs.save(PPTX)
print(f"Güncellendi: {PPTX}  ({len(prs.slides)} slayt)")
