"""
Build microgpt_explained.pptx
A self-contained slide deck that teaches everything in microgpt.py from scratch.
Run once after training: python create_presentation.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

DARK   = RGBColor(0x1a, 0x1a, 0x2e)   # near-black navy
ACCENT = RGBColor(0x16, 0x21, 0x3e)   # dark blue panel
GOLD   = RGBColor(0xe9, 0x4f, 0x37)   # warm red-orange for highlights
WHITE  = RGBColor(0xff, 0xff, 0xff)
LIGHT  = RGBColor(0xec, 0xf0, 0xf1)   # off-white text
MUTED  = RGBColor(0x95, 0xa5, 0xa6)   # grey subtitles
GREEN  = RGBColor(0x27, 0xae, 0x60)
BLUE   = RGBColor(0x29, 0x80, 0xb9)

HERE = os.path.dirname(os.path.abspath(__file__))

def img(name):
    return os.path.join(HERE, name)

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def bg(slide, color=DARK):
    """Fill slide background with a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height, color):
    """Colored rectangle (no text)."""
    sh = slide.shapes.add_shape(1, left, top, width, height)   # MSO_SHAPE_TYPE.RECTANGLE=1
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def label(slide, text, left, top, width, height,
          size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    """Add a text box."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def bullet_box(slide, items, left, top, width, height,
               size=16, title=None, title_size=20, title_color=GOLD):
    """
    Multi-line bullet list. items is a list of strings.
    Prefix each item with • automatically.
    """
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.size  = Pt(title_size)
        r.font.bold  = True
        r.font.color.rgb = title_color

    for i, item in enumerate(items):
        p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = f"  •  {item}"
        r.font.size  = Pt(size)
        r.font.color.rgb = LIGHT

def add_image(slide, path, left, top, width=None, height=None):
    if not os.path.exists(path):
        return
    if width and height:
        slide.shapes.add_picture(path, left, top, width, height)
    elif width:
        slide.shapes.add_picture(path, left, top, width=width)
    elif height:
        slide.shapes.add_picture(path, left, top, height=height)
    else:
        slide.shapes.add_picture(path, left, top)

def divider(slide, top, color=GOLD):
    box(slide, Inches(0.5), top, Inches(12.33), Inches(0.04), color)

def code_box(slide, code, left, top, width, height, size=13):
    """Monospaced code block on a dark panel."""
    box(slide, left, top, width, height, ACCENT)
    txb = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), height - Inches(0.2))
    tf = txb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code.strip().split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = RGBColor(0xa8, 0xd8, 0xea)   # light blue mono
        r.font.name = 'Courier New'

# ---------------------------------------------------------------------------
# Build slides
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── 1. TITLE ────────────────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
box(sl, 0, Inches(2.5), W, Inches(2.6), ACCENT)
label(sl, "MicroGPT", Inches(0.8), Inches(2.65), Inches(11), Inches(1.1),
      size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
label(sl, "How a GPT learns to generate names — from scratch, in pure Python",
      Inches(0.8), Inches(3.7), Inches(11), Inches(0.7),
      size=22, color=LIGHT, align=PP_ALIGN.CENTER)
label(sl, "Based on @karpathy · microgpt.py",
      Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
      size=14, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# ── 2. BIG PICTURE ─────────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "The Big Picture", Inches(0.6), Inches(0.3), Inches(11), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(1.05))

steps = [
    ("1  Dataset",       "32,000 first names.  Each name is one training example."),
    ("2  Tokenizer",     "Split each name into characters.  Map each character to an integer ID."),
    ("3  Model",         "A tiny Transformer (GPT) that predicts the next character, one at a time."),
    ("4  Autograd",      "A hand-built automatic-differentiation engine tracks every computation."),
    ("5  Training",      "Feed names in, measure how wrong the predictions are (loss), nudge weights."),
    ("6  Inference",     "Start from a blank slate, let the model hallucinate new names character by character."),
]
for i, (title, desc) in enumerate(steps):
    top = Inches(1.25) + i * Inches(0.97)
    box(sl, Inches(0.5), top, Inches(2.4), Inches(0.78), ACCENT)
    label(sl, title, Inches(0.55), top + Inches(0.12), Inches(2.3), Inches(0.6),
          size=15, bold=True, color=GOLD)
    label(sl, desc, Inches(3.1), top + Inches(0.12), Inches(9.6), Inches(0.65),
          size=15, color=LIGHT)

# ── 3. DATASET ──────────────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Step 1 — Dataset & Tokenizer", Inches(0.6), Inches(0.3), Inches(11), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(1.05))

bullet_box(sl, [
    "input.txt  contains ~32,000 English first names, one per line.",
    "Each name is a document.  Names are shuffled randomly before training.",
    "Vocabulary = every unique character that appears  +  one special <BOS> token.",
    "<BOS>  (Beginning Of Sequence) marks both the START and END of a name.",
    "Example: 'emma'  →  [<BOS>, e, m, m, a, <BOS>]",
    "Tokenizer maps characters to integers 0…N and back.  No subword magic needed.",
], Inches(0.6), Inches(1.2), Inches(12.1), Inches(3.0), size=17)

code_box(sl, """
uchars    = sorted(set(''.join(docs)))    # e.g. ['a','b','c',...,'z']
BOS       = len(uchars)                   # e.g. 26  ← special token id
vocab_size = len(uchars) + 1             # 27 total token ids

# Tokenize 'emma':
tokens = [BOS] + [uchars.index(ch) for ch in 'emma'] + [BOS]
# → [26, 4, 12, 12, 0, 26]
""", Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.5), size=14)

# ── 4. AUTOGRAD ─────────────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Step 2 — Autograd: Teaching the Computer to Differentiate",
      Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.7),
      size=30, bold=True, color=GOLD)
divider(sl, Inches(1.05))

bullet_box(sl, [
    "Every number in the forward pass is wrapped in a  Value  object.",
    "Value stores both the number (.data) and its gradient (.grad).",
    "When you do  a + b  or  a * b,  a new Value records who its parents are.",
    "This builds a computation graph — a trail of breadcrumbs back to every weight.",
    ".backward()  walks that graph in reverse, applying the chain rule automatically.",
    "Result: every weight knows exactly how much it contributed to the loss.",
], Inches(0.6), Inches(1.2), Inches(6.2), Inches(3.4), size=16)

code_box(sl, """
class Value:
    def __init__(self, data):
        self.data = data   # the number
        self.grad = 0      # d(loss)/d(self), filled by backward()

    def __mul__(self, other):
        # local grads: d(a*b)/da = b,  d(a*b)/db = a
        return Value(self.data * other.data,
                     children=(self, other),
                     local_grads=(other.data, self.data))

a = Value(3.0)
b = Value(4.0)
c = a * b          # c.data = 12.0
c.backward()       # a.grad = 4.0,  b.grad = 3.0
""", Inches(6.9), Inches(1.2), Inches(5.9), Inches(4.8), size=12)

label(sl, "Chain rule in one sentence:",
      Inches(0.6), Inches(4.7), Inches(6.2), Inches(0.4),
      size=14, bold=True, color=GOLD)
label(sl, "grad(child) += local_grad × grad(parent)",
      Inches(0.6), Inches(5.1), Inches(6.2), Inches(0.5),
      size=16, color=GREEN, bold=True)
label(sl, "Repeat for every node, from loss back to weights.",
      Inches(0.6), Inches(5.6), Inches(6.2), Inches(0.5),
      size=14, color=LIGHT)

# ── 5. HYPERPARAMETERS ──────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Step 3 — Hyperparameters", Inches(0.6), Inches(0.2), Inches(7), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))
label(sl, "These are the knobs you set before training begins.\nThe model cannot learn them — you choose them.",
      Inches(0.6), Inches(1.05), Inches(6), Inches(0.9), size=16, color=LIGHT)
add_image(sl, img('hyperparams.png'), Inches(0.3), Inches(1.8), height=Inches(5.4))

# ── 6. MODEL WEIGHTS ────────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Step 4 — Model Weights (What the Model Knows)",
      Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=34, bold=True, color=GOLD)
divider(sl, Inches(0.95))

bullet_box(sl, [
    "wte  (vocab × n_embd)  — Token Embedding Table.  One learned vector per character.",
    "wpe  (block_size × n_embd)  — Position Embedding Table.  One vector per position slot.",
    "attn_wq / wk / wv / wo  — Query, Key, Value, Output projections inside each attention head.",
    "mlp_fc1 / fc2  — Two linear layers forming the feed-forward block.",
    "lm_head  — Projects the final hidden state back to vocabulary size → logits.",
    "All weights start as small random Gaussians (std=0.08).  Training shapes them.",
], Inches(0.6), Inches(1.1), Inches(7.0), Inches(3.2), size=15)

add_image(sl, img('initial_weights.png'), Inches(7.5), Inches(1.0), height=Inches(5.6))
label(sl, "Initial weights — pure noise.\nEach cell = one weight value.\nRed = negative, Blue = positive.",
      Inches(7.5), Inches(6.6), Inches(5.5), Inches(0.7), size=11, color=MUTED, italic=True)

# ── 7. ARCHITECTURE OVERVIEW ────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Step 5 — Model Architecture (The Transformer)",
      Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=34, bold=True, color=GOLD)
divider(sl, Inches(0.95))

blocks = [
    ("Token + Position Embedding",
     "Look up the character's vector (wte) and add the position's vector (wpe).\n"
     "Result: a 16-dim vector that encodes both WHAT the character is and WHERE it sits."),
    ("RMSNorm",
     "Normalise the vector so its root-mean-square = 1.\n"
     "Keeps numbers well-behaved throughout the network."),
    ("Multi-Head Self-Attention",
     "Each position looks at all previous positions and decides what to borrow.\n"
     "4 heads do this in parallel, each attending to different patterns."),
    ("Residual Connection",
     "Add the attention output back to the input (skip connection).\n"
     "Lets gradients flow directly to early layers — makes training stable."),
    ("MLP Block",
     "Two linear layers with ReLU in between — width expands 4×, then contracts.\n"
     "Stores and transforms patterns the attention found."),
    ("LM Head → Logits → Softmax",
     "Project to vocab_size numbers (logits), convert to probabilities.\n"
     "The highest-probability token is the model's best guess for what comes next."),
]
for i, (title, desc) in enumerate(blocks):
    top = Inches(1.1) + i * Inches(1.03)
    box(sl, Inches(0.4), top, Inches(0.08), Inches(0.78), GOLD)
    label(sl, title, Inches(0.6), top, Inches(4.2), Inches(0.38),
          size=14, bold=True, color=GOLD)
    label(sl, desc, Inches(0.6), top + Inches(0.37), Inches(12.2), Inches(0.55),
          size=13, color=LIGHT)

# ── 8. ATTENTION DEEP DIVE ──────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Inside Attention — Q, K, V",
      Inches(0.6), Inches(0.2), Inches(9), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))

label(sl, "Think of it like a search engine inside the model:",
      Inches(0.6), Inches(1.1), Inches(12), Inches(0.45), size=18, color=LIGHT, bold=True)

qkv = [
    ("Q  Query", GOLD,
     "\"What am I looking for?\"\nThe current token asks a question."),
    ("K  Key",   BLUE,
     "\"What do I contain?\"\nEvery past token advertises itself."),
    ("V  Value", GREEN,
     "\"What will I give you?\"\nThe actual information transferred if attention is high."),
]
for i, (name, col, desc) in enumerate(qkv):
    left = Inches(0.5) + i * Inches(4.25)
    box(sl, left, Inches(1.7), Inches(3.9), Inches(1.8), ACCENT)
    label(sl, name, left + Inches(0.15), Inches(1.8), Inches(3.6), Inches(0.5),
          size=20, bold=True, color=col)
    label(sl, desc, left + Inches(0.15), Inches(2.3), Inches(3.6), Inches(1.1),
          size=14, color=LIGHT)

label(sl, "Score  =  (Q · K) / √head_dim        →       softmax       →       weighted sum of V",
      Inches(0.6), Inches(3.65), Inches(12), Inches(0.5),
      size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

bullet_box(sl, [
    "The dot product Q·K measures how well a query matches each key — higher = more attention.",
    "Dividing by √head_dim prevents the dot products from getting too large before softmax.",
    "Softmax converts scores to probabilities that sum to 1.",
    "Each position can only attend to itself and earlier positions (causal masking) — no peeking at the future.",
    "Running 4 heads in parallel lets the model attend to different things simultaneously.",
], Inches(0.6), Inches(4.25), Inches(12.1), Inches(2.8), size=15)

# ── 9. ATTENTION HEATMAP ────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Attention in Action", Inches(0.6), Inches(0.2), Inches(7), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))

label(sl, "The heatmap below shows the attention weight matrix for one word across training.\n"
          "Row = query position (the token doing the attending).\n"
          "Column = key position (the token being attended to).\n"
          "Brighter = more attention.",
      Inches(0.6), Inches(1.1), Inches(6.2), Inches(1.5), size=15, color=LIGHT)
label(sl, "What to notice:",
      Inches(0.6), Inches(2.65), Inches(6.2), Inches(0.35), size=15, bold=True, color=GOLD)
bullet_box(sl, [
    "Step 1: uniform blur — the model attends everywhere equally (knows nothing).",
    "During training: the lower-left triangle sharpens — causal structure emerges.",
    "Late training: specific cells brighten — the model learns which characters predict which.",
    "The diagonal is almost always bright: every character attends strongly to itself.",
], Inches(0.6), Inches(3.0), Inches(6.2), Inches(2.8), size=14)
add_image(sl, img('attention_animated.gif'), Inches(6.8), Inches(1.0), height=Inches(6.0))

# ── 10. TRAINING LOOP ───────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Step 6 — Training Loop", Inches(0.6), Inches(0.2), Inches(9), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))

phases = [
    ("① Forward Pass",
     "Feed a tokenized name into the model.\nGet a probability distribution over the next character at each position."),
    ("② Loss",
     "Measure how surprised the model is by the correct next character.\n"
     "Loss = -log(probability assigned to the correct token).  Lower = better."),
    ("③ Backward Pass",
     "Call .backward() on the loss.\n"
     "Autograd walks the computation graph and fills .grad for every weight."),
    ("④ Adam Update",
     "Adjust each weight by a small step in the direction that reduces loss.\n"
     "Adam keeps a running average of gradients (momentum) and scales each step individually."),
    ("⑤ Zero Gradients",
     "Reset .grad = 0 for all weights before the next step.\n"
     "Otherwise gradients accumulate across steps."),
]
for i, (title, desc) in enumerate(phases):
    top = Inches(1.1) + i * Inches(1.15)
    box(sl, Inches(0.4), top, Inches(2.6), Inches(0.9), ACCENT)
    label(sl, title, Inches(0.5), top + Inches(0.15), Inches(2.4), Inches(0.6),
          size=15, bold=True, color=GOLD)
    label(sl, desc, Inches(3.2), top + Inches(0.08), Inches(9.6), Inches(0.85),
          size=14, color=LIGHT)

# ── 11. ADAM OPTIMIZER ──────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "The Adam Optimizer", Inches(0.6), Inches(0.2), Inches(9), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))

label(sl, "Plain gradient descent: weight -= lr × grad\n"
          "Problem: some weights need big steps, others tiny. lr can't be both.\n"
          "Adam solves this by keeping a separate adaptive step size per weight.",
      Inches(0.6), Inches(1.1), Inches(12), Inches(1.1), size=16, color=LIGHT)

code_box(sl, """
# For each weight i at each step t:
m[i] = beta1 * m[i] + (1 - beta1) * grad      # momentum  (smoothed gradient)
v[i] = beta2 * v[i] + (1 - beta2) * grad**2   # RMS       (smoothed squared gradient)

m_hat = m[i] / (1 - beta1**t)   # bias-corrected (important early in training)
v_hat = v[i] / (1 - beta2**t)

weight -= lr * m_hat / (sqrt(v_hat) + eps)     # adaptive step per weight
""", Inches(0.6), Inches(2.3), Inches(12.1), Inches(2.8), size=14)

bullet_box(sl, [
    "m  is the momentum buffer — it smooths out noisy gradients.",
    "v  is the RMS buffer — weights with large gradients get smaller steps automatically.",
    "beta1=0.85 and beta2=0.99 control how fast the buffers update.",
    "eps=1e-8 prevents division by zero.",
    "lr decays linearly to zero so the model settles rather than bouncing around.",
], Inches(0.6), Inches(5.2), Inches(12.1), Inches(2.0), size=15)

# ── 12. TRAINING CURVES ─────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Watching the Model Learn", Inches(0.6), Inches(0.2), Inches(9), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))

label(sl, "Top: training loss over 1,000 steps.   Bottom: learning rate (linear decay).",
      Inches(0.6), Inches(1.1), Inches(12), Inches(0.4), size=15, color=LIGHT)
add_image(sl, img('training_curves_animated.gif'), Inches(0.4), Inches(1.55), height=Inches(4.5))

bullet_box(sl, [
    "Loss starts high (~3) — the model is essentially guessing randomly.",
    "Loss drops quickly at first, then slows — easy patterns learned first.",
    "A good final loss for this task is around 1.5–2.0.",
    "Learning rate decays to zero so updates shrink as training ends.",
], Inches(7.5), Inches(2.0), Inches(5.5), Inches(3.5), size=15)

label(sl, "Loss = how surprised the model is.\n"
          "Lower = more confident = better predictions.",
      Inches(7.5), Inches(5.5), Inches(5.5), Inches(0.9),
      size=14, color=MUTED, italic=True)

# ── 13. WEIGHT DISTRIBUTION ─────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Weight Distribution — Noise → Signal",
      Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=34, bold=True, color=GOLD)
divider(sl, Inches(0.95))
add_image(sl, img('weight_dist_animated.gif'), Inches(0.3), Inches(1.1), height=Inches(5.5))
bullet_box(sl, [
    "Step 1: wide bell curve — weights drawn from Gaussian(0, 0.08).",
    "Training: the distribution narrows and shifts as useful weights grow, useless ones shrink.",
    "Mean stays near zero; the spread (σ) tells you how 'confident' the model is.",
    "A healthy distribution is not too wide (exploding) and not collapsed to zero (dead).",
    "Compare initial_weights.png vs final_weights.png — you can see structure form.",
], Inches(7.5), Inches(1.5), Inches(5.6), Inches(4.0), size=15)

# ── 14. INITIAL vs FINAL WEIGHTS ────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Initial vs Final Weights", Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))
label(sl, "Before training", Inches(0.8), Inches(1.1), Inches(5.8), Inches(0.4),
      size=18, bold=True, color=MUTED)
label(sl, "After training", Inches(7.1), Inches(1.1), Inches(5.8), Inches(0.4),
      size=18, bold=True, color=GREEN)
add_image(sl, img('initial_weights.png'), Inches(0.3), Inches(1.5), width=Inches(6.3))
add_image(sl, img('final_weights.png'),   Inches(6.8), Inches(1.5), width=Inches(6.3))
label(sl, "Random noise — no structure.",
      Inches(0.5), Inches(6.4), Inches(6.0), Inches(0.5), size=13, color=MUTED, italic=True)
label(sl, "Structured patterns — each matrix has learned a role.",
      Inches(7.0), Inches(6.4), Inches(6.0), Inches(0.5), size=13, color=GREEN, italic=True)

# ── 15. EMBEDDING PCA ───────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "What Did the Model Learn? — Embedding Space",
      Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=32, bold=True, color=GOLD)
divider(sl, Inches(0.95))

label(sl, "Each character has a 16-dimensional learned vector (wte).\n"
          "PCA collapses those 16 dimensions to 2 so we can plot them.\n"
          "Watch clusters form as training progresses.",
      Inches(0.6), Inches(1.1), Inches(6.0), Inches(1.2), size=15, color=LIGHT)
bullet_box(sl, [
    "Red dots = vowels (a, e, i, o, u)",
    "Blue dots = consonants",
    "Green star = <BOS> token",
    "Step 1: scattered randomly — the model knows nothing.",
    "Mid-training: vowels and consonants pull apart.",
    "Late training: similar sounds cluster — b/p, m/n, a/e, etc.",
    "<BOS> drifts far from all letters — it has a unique role (start + end).",
    "This is the model's internal 'dictionary of meanings' becoming real.",
], Inches(0.6), Inches(2.35), Inches(5.8), Inches(4.4), size=14)
add_image(sl, img('embedding_pca_animated.gif'), Inches(6.6), Inches(1.0), height=Inches(6.1))

# ── 16. FINAL PARAMETER DISTRIBUTION ───────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Final Parameter Distribution",
      Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))
add_image(sl, img('final_parameter_distribution.png'), Inches(0.3), Inches(1.1), height=Inches(5.0))
bullet_box(sl, [
    "Left: final weight values — narrower than the initial Gaussian.",
    "Right: final gradient values — near zero means the model has converged.",
    "Large gradients late in training = the model is still changing a lot (not converged).",
    "The red dashed line marks the mean — close to 0 means no systematic bias.",
], Inches(7.8), Inches(2.0), Inches(5.3), Inches(3.5), size=15)

# ── 17. INFERENCE ───────────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Step 7 — Inference: Generating Names",
      Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=34, bold=True, color=GOLD)
divider(sl, Inches(0.95))

steps_inf = [
    ("Start", "Feed <BOS> as the first token.  The model has seen nothing yet."),
    ("Predict", "Model outputs a probability distribution over all 27 tokens."),
    ("Sample",  "Pick the next token randomly, weighted by those probabilities."),
    ("Append",  "That token becomes the next input.  The context window grows."),
    ("Repeat",  "Keep predicting until the model outputs <BOS> again (= end of name)."),
]
for i, (title, desc) in enumerate(steps_inf):
    top = Inches(1.2) + i * Inches(1.05)
    box(sl, Inches(0.4), top, Inches(1.6), Inches(0.8), ACCENT)
    label(sl, title, Inches(0.45), top + Inches(0.15), Inches(1.5), Inches(0.5),
          size=16, bold=True, color=GOLD)
    label(sl, desc, Inches(2.2), top + Inches(0.1), Inches(10.5), Inches(0.7),
          size=16, color=LIGHT)

label(sl, "Temperature  controls randomness:",
      Inches(0.6), Inches(6.45), Inches(12), Inches(0.35),
      size=15, bold=True, color=GOLD)
label(sl, "Low temp (0.1) = model picks the most likely character every time → repetitive but safe.   "
          "High temp (2.0) = wild guesses → creative but nonsense.   0.5 is the sweet spot here.",
      Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
      size=14, color=LIGHT)

# ── 18. GENERATION ANIMATION ────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Generation — Character by Character",
      Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))

label(sl, "The bar chart shows the model's probability for each candidate next character.\n"
          "Gold bar = the character that was actually sampled.",
      Inches(0.6), Inches(1.1), Inches(6.0), Inches(0.9), size=15, color=LIGHT)
bullet_box(sl, [
    "Early steps: probabilities spread across many characters — still uncertain.",
    "After a few characters: the distribution sharpens — the name's pattern is clear.",
    "e.g. after 'Em' the model strongly prefers 'm', 'i', 'a' — all valid continuations.",
    "The model never looks up names — it learned the patterns from scratch.",
], Inches(0.6), Inches(2.1), Inches(6.0), Inches(3.0), size=14)
add_image(sl, img('generation_animated.gif'), Inches(6.5), Inches(1.0), height=Inches(6.1))

# ── 19. KEY EQUATIONS ───────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Key Equations — One Slide",
      Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))

eqs = [
    ("Embedding",      "x  =  wte[token_id]  +  wpe[pos_id]"),
    ("RMSNorm",        "x̂  =  x / √(mean(x²) + ε)"),
    ("Attention score","score(q, k)  =  (q · k) / √head_dim"),
    ("Softmax",        "softmax(zᵢ)  =  exp(zᵢ) / Σ exp(zⱼ)"),
    ("Attention out",  "out  =  Σ softmax(scores)ₜ · vₜ"),
    ("Cross-entropy",  "loss  =  −log p(correct token)"),
    ("Adam step",      "w  −=  lr · m̂ / (√v̂ + ε)"),
]
for i, (name, eq) in enumerate(eqs):
    top = Inches(1.1) + i * Inches(0.83)
    label(sl, name, Inches(0.6), top, Inches(2.8), Inches(0.5),
          size=15, bold=True, color=GOLD)
    box(sl, Inches(3.5), top, Inches(9.4), Inches(0.65), ACCENT)
    label(sl, eq, Inches(3.65), top + Inches(0.06), Inches(9.1), Inches(0.55),
          size=16, color=GREEN, bold=True)

# ── 20. SUMMARY ─────────────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
box(sl, 0, Inches(2.8), W, Inches(2.0), ACCENT)
label(sl, "Summary", Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))

takeaways = [
    "A GPT is just matrix multiplications, softmax, and a loop.",
    "Autograd makes training possible — no manual derivatives needed.",
    "Attention lets every position gather information from every past position.",
    "Embeddings are the model's learned alphabet of meaning — PCA makes them visible.",
    "Adam adapts the step size per weight — far better than plain gradient descent.",
    "Temperature is the single knob between 'precise' and 'creative' generation.",
    "1,000 steps and ~7,000 parameters is enough to hallucinate plausible names.",
]
for i, t in enumerate(takeaways):
    top = Inches(1.1) + i * Inches(0.82)
    label(sl, f"{'★' if i < 3 else '·'}  {t}",
          Inches(0.7), top, Inches(12), Inches(0.6),
          size=16, color=WHITE if i < 3 else LIGHT,
          bold=(i < 3))

label(sl, "microgpt.py  ·  ~380 lines  ·  zero dependencies  ·  complete GPT",
      Inches(0.6), Inches(6.8), Inches(12), Inches(0.45),
      size=13, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------

out = os.path.join(HERE, 'microgpt_explained.pptx')
prs.save(out)
print(f"Saved  {out}")
print(f"Slides: {len(prs.slides)}")
