"""
microgpt_aciklamali.pptx dosyasına 8. slayt olarak mimari diyagramını ekler.
(7. slayt = Transformer mimarisi metni, 9. slayt = Q/K/V detayı)
Mevcut slaytları bozmadan araya yerleştirir.
"""

import os
import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE  = os.path.dirname(os.path.abspath(__file__))
PPTX  = os.path.join(HERE, 'microgpt_aciklamali.pptx')
IMG   = os.path.join(HERE, 'architecture_tr.png')

W, H   = Inches(13.33), Inches(7.5)
DARK   = RGBColor(0x1a, 0x1a, 0x2e)
GOLD   = RGBColor(0xe9, 0x4f, 0x37)
LIGHT  = RGBColor(0xec, 0xf0, 0xf1)
MUTED  = RGBColor(0x95, 0xa5, 0xa6)

prs = Presentation(PPTX)

# ── Yeni boş slayt oluştur ────────────────────────────────────────────────────
blank_layout = prs.slide_layouts[6]
new_slide     = prs.slides.add_slide(blank_layout)

# Arka plan rengi
fill = new_slide.background.fill
fill.solid()
fill.fore_color.rgb = DARK

def box(slide, left, top, w, h, color):
    sh = slide.shapes.add_shape(1, left, top, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()

def label(slide, text, left, top, w, h,
          size=18, bold=False, color=LIGHT, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(left, top, w, h)
    tf  = txb.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color

def divider(slide, top):
    box(slide, Inches(0.5), top, Inches(12.33), Inches(0.04), GOLD)

# ── Başlık ────────────────────────────────────────────────────────────────────
label(new_slide,
      'Genel Mimari — Tek Bakışta GPT',
      Inches(0.6), Inches(0.18), Inches(12), Inches(0.65),
      size=34, bold=True, color=GOLD)
divider(new_slide, Inches(0.88))

# ── Ana diyagram görseli ───────────────────────────────────────────────────────
new_slide.shapes.add_picture(IMG, Inches(0.25), Inches(0.95), height=Inches(5.8))

# ── Sağ taraf açıklama kutuları ───────────────────────────────────────────────
notlar = [
    (GOLD,                  '① Girdi',
     'Her harf bir token ID\'sine dönüşür. Konum bilgisi ayrı eklenir.'),
    (RGBColor(0x1a,0x6b,0x9a), '② Gömme',
     'Token + pozisyon vektörleri toplanır → 16 boyutlu x vektörü.'),
    (RGBColor(0x6c,0x3a,0x7a), '③ RMSNorm',
     'Sayıları dengede tutar; karekök-ortalama-kare = 1 yapılır.'),
    (RGBColor(0x1a,0x6b,0x3a), '④ Dikkat',
     'Q·K skoru hesaplanır, softmax ile V ağırlıklı toplanır. 4 kafa paralel çalışır.'),
    (RGBColor(0x7a,0x4a,0x00), '⑤ MLP',
     'FC1 ile genişlet (×4), ReLU uygula, FC2 ile daralt. Bilgi depolanır.'),
    (RGBColor(0x7a,0x1a,0x1a), '⑥ LM Kafası',
     'Son x vektörü vocab_size logitine dönüştürülür.'),
    (RGBColor(0x1a,0x5a,0x1a), '⑦ Softmax',
     'Logitler olasılığa çevrilir. En yüksek olasılıklı harf seçilir.'),
]

for i, (renk, baslik, aciklama) in enumerate(notlar):
    top = Inches(0.98) + i * Inches(0.80)
    sh = new_slide.shapes.add_shape(1, Inches(10.1), top, Inches(3.0), Inches(0.70))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x16, 0x21, 0x3e)
    sh.line.color.rgb = renk; sh.line.width = Pt(1.2)

    txb = new_slide.shapes.add_textbox(Inches(10.18), top + Inches(0.03),
                                        Inches(2.85), Inches(0.65))
    tf  = txb.text_frame; tf.word_wrap = True

    p  = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = baslik + '  '
    r1.font.size = Pt(11); r1.font.bold = True; r1.font.color.rgb = renk

    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = aciklama
    r2.font.size = Pt(9); r2.font.color.rgb = LIGHT

# ── Artık bağlantı notu ───────────────────────────────────────────────────────
label(new_slide,
      '↺  Her blokta Artık Bağlantı: x = x_blok_çıktısı + x_giriş',
      Inches(0.5), Inches(6.88), Inches(9.5), Inches(0.42),
      size=11, color=MUTED, italic=True)

# ── Yeni slaytı doğru konuma taşı (index 7 = 8. slayt) ──────────────────────
# python-pptx'te slayt sırası XML düzeyinde değiştirilir
xml_slides = prs.slides._sldIdLst
# son eklenen slayt şu an en sonda; onu index 7'ye taşı
slides_list = list(xml_slides)
new_el      = slides_list[-1]          # yeni eklenen
xml_slides.remove(new_el)
# index 7'ye yerleştir (0-tabanlı → 8. slayt)
insert_pos = min(7, len(slides_list) - 1)
ref_el     = slides_list[insert_pos]   # şu anki 8. slayt
xml_slides.insert(insert_pos, new_el)

prs.save(PPTX)
print(f"Güncellendi: {PPTX}  ({len(prs.slides)} slayt)")
