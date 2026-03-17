"""
Build microgpt_aciklamali.pptx  (Türkçe versiyon)
MicroGPT'yi sıfırdan anlatan slayt destesi.
Çalıştırmak için: python create_presentation_tr.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Yardımcı fonksiyonlar
# ---------------------------------------------------------------------------

W, H = Inches(13.33), Inches(7.5)   # 16:9 geniş ekran

DARK   = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT = RGBColor(0x16, 0x21, 0x3e)
GOLD   = RGBColor(0xe9, 0x4f, 0x37)
WHITE  = RGBColor(0xff, 0xff, 0xff)
LIGHT  = RGBColor(0xec, 0xf0, 0xf1)
MUTED  = RGBColor(0x95, 0xa5, 0xa6)
GREEN  = RGBColor(0x27, 0xae, 0x60)
BLUE   = RGBColor(0x29, 0x80, 0xb9)

HERE = os.path.dirname(os.path.abspath(__file__))

def img(name):
    return os.path.join(HERE, name)

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height, color):
    sh = slide.shapes.add_shape(1, left, top, width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def label(slide, text, left, top, width, height,
          size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def bullet_box(slide, items, left, top, width, height,
               size=16, title=None, title_size=20, title_color=GOLD):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.size  = Pt(title_size)
        r.font.bold  = True
        r.font.color.rgb = title_color
    for i, item in enumerate(items):
        p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = f"  •  {item}"
        r.font.size  = Pt(size)
        r.font.color.rgb = LIGHT

def add_image(slide, path, left, top, width=None, height=None):
    if not os.path.exists(path):
        return
    if width and height:
        slide.shapes.add_picture(path, left, top, width, height)
    elif width:
        slide.shapes.add_picture(path, left, top, width=width)
    elif height:
        slide.shapes.add_picture(path, left, top, height=height)
    else:
        slide.shapes.add_picture(path, left, top)

def divider(slide, top, color=GOLD):
    box(slide, Inches(0.5), top, Inches(12.33), Inches(0.04), color)

def code_box(slide, code, left, top, width, height, size=13):
    box(slide, left, top, width, height, ACCENT)
    txb = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), height - Inches(0.2))
    tf = txb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code.strip().split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = RGBColor(0xa8, 0xd8, 0xea)
        r.font.name = 'Courier New'

# ---------------------------------------------------------------------------
# Slaytlar
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── 1. BAŞLIK ───────────────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
box(sl, 0, Inches(2.5), W, Inches(2.6), ACCENT)
label(sl, "MicroGPT", Inches(0.8), Inches(2.65), Inches(11), Inches(1.1),
      size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
label(sl, "Sıfırdan, saf Python ile bir GPT nasıl öğrenir ve isim üretir?",
      Inches(0.8), Inches(3.7), Inches(11), Inches(0.7),
      size=22, color=LIGHT, align=PP_ALIGN.CENTER)
label(sl, "@karpathy · microgpt.py  ·  Türkçe Anlatım",
      Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
      size=14, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# ── 2. BÜYÜK RESİM ──────────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Büyük Resim", Inches(0.6), Inches(0.3), Inches(11), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(1.05))

steps = [
    ("1  Veri Seti",       "32.000 İngilizce isim.  Her isim tek bir eğitim örneği."),
    ("2  Tokenizer",       "Her ismi harflere böl.  Her harfi bir tam sayı kimliğine (ID) çevir."),
    ("3  Model",           "Küçük bir Transformer (GPT) — bir sonraki harfi tahmin eder."),
    ("4  Otomatik Türev",  "El yapımı bir türev motoru her hesabı takip eder (Autograd)."),
    ("5  Eğitim",          "İsimleri ver, tahmin hatası (kayıp) hesapla, ağırlıkları küçük adımlarla güncelle."),
    ("6  Çıkarım",         "Boş bir başlangıçtan, model harf harf yeni isimler 'hayal eder'."),
]
for i, (baslik, aciklama) in enumerate(steps):
    top = Inches(1.25) + i * Inches(0.97)
    box(sl, Inches(0.5), top, Inches(2.5), Inches(0.78), ACCENT)
    label(sl, baslik, Inches(0.55), top + Inches(0.12), Inches(2.4), Inches(0.6),
          size=15, bold=True, color=GOLD)
    label(sl, aciklama, Inches(3.1), top + Inches(0.12), Inches(9.6), Inches(0.65),
          size=15, color=LIGHT)

# ── 3. VERİ SETİ & TOKENİZER ────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Adım 1 — Veri Seti ve Tokenizer", Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(1.05))

bullet_box(sl, [
    "input.txt dosyası ~32.000 İngilizce isim içerir, her satırda bir isim.",
    "Her isim bir 'belge' sayılır.  Eğitimden önce sıraları karıştırılır.",
    "Kelime dağarcığı (vocab): veri setinde geçen tüm benzersiz harfler + özel <BOS> tokeni.",
    "<BOS> (Dizi Başlangıcı) hem ismin BAŞINI hem de SONUNU işaret eder.",
    "Örnek: 'ali'  →  [<BOS>, a, l, i, <BOS>]",
    "Tokenizer harfleri 0…N tam sayılarına, sayıları tekrar harflere çevirir.",
], Inches(0.6), Inches(1.2), Inches(12.1), Inches(3.0), size=17)

code_box(sl, """
uchars    = sorted(set(''.join(docs)))    # örn. ['a','b','c',...,'z']
BOS       = len(uchars)                   # örn. 26  ← özel token id'si
vocab_size = len(uchars) + 1             # 27 benzersiz token

# 'ali' kelimesini tokenize et:
tokens = [BOS] + [uchars.index(ch) for ch in 'ali'] + [BOS]
# → [26, 0, 11, 8, 26]
""", Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.5), size=14)

# ── 4. OTOMATİK TÜREV (AUTOGRAD) ────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Adım 2 — Autograd: Bilgisayara Türev Aldırma",
      Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.7),
      size=30, bold=True, color=GOLD)
divider(sl, Inches(1.05))

bullet_box(sl, [
    "İleri geçişteki her sayı bir Value nesnesine sarılır.",
    "Value hem sayıyı (.data) hem de gradyanı (.grad) saklar.",
    "a + b veya a * b yapıldığında yeni Value hangi ata-düğümlerden geldiğini kaydeder.",
    "Bu, tüm ağırlıklara giden bir hesap grafiği (yol haritası) oluşturur.",
    ".backward() bu grafiği ters yönde dolaşır ve zincir kuralını otomatik uygular.",
    "Sonuç: her ağırlık kaybı ne kadar etkilediğini tam olarak bilir.",
], Inches(0.6), Inches(1.2), Inches(6.2), Inches(3.4), size=16)

code_box(sl, """
class Value:
    def __init__(self, data):
        self.data = data   # asıl sayı
        self.grad = 0      # d(kayip)/d(self), backward() ile dolar

    def __mul__(self, other):
        # yerel türevler: d(a*b)/da = b,  d(a*b)/db = a
        return Value(self.data * other.data,
                     children=(self, other),
                     local_grads=(other.data, self.data))

a = Value(3.0)
b = Value(4.0)
c = a * b          # c.data = 12.0
c.backward()       # a.grad = 4.0,  b.grad = 3.0
""", Inches(6.9), Inches(1.2), Inches(5.9), Inches(4.8), size=12)

label(sl, "Zincir kuralı tek cümlede:",
      Inches(0.6), Inches(4.7), Inches(6.2), Inches(0.4),
      size=14, bold=True, color=GOLD)
label(sl, "grad(çocuk) += yerel_grad × grad(ebeveyn)",
      Inches(0.6), Inches(5.1), Inches(6.2), Inches(0.5),
      size=16, color=GREEN, bold=True)
label(sl, "Kayıptan başlayarak tüm ağırlıklara kadar her düğümde tekrarla.",
      Inches(0.6), Inches(5.6), Inches(6.2), Inches(0.5),
      size=14, color=LIGHT)

# ── 5. HİPERPARAMETRELER ────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Adım 3 — Hiperparametreler", Inches(0.6), Inches(0.2), Inches(7), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))
label(sl, "Eğitim başlamadan önce senin belirlediğin ayarlar.\nModel bunları öğrenemez — sen seçersin.",
      Inches(0.6), Inches(1.05), Inches(6), Inches(0.9), size=16, color=LIGHT)
add_image(sl, img('hyperparams.png'), Inches(0.3), Inches(1.8), height=Inches(5.4))

# ── 6. MODEL AĞIRLIKLARI ────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Adım 4 — Model Ağırlıkları (Modelin Bilgisi)",
      Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=34, bold=True, color=GOLD)
divider(sl, Inches(0.95))

bullet_box(sl, [
    "wte  (vocab × n_embd)  — Token Gömme Tablosu.  Her harf için bir öğrenilmiş vektör.",
    "wpe  (block_size × n_embd)  — Pozisyon Gömme Tablosu.  Her konum için bir vektör.",
    "attn_wq / wk / wv / wo  — Dikkat başlıkları içindeki Sorgu, Anahtar, Değer, Çıktı katmanları.",
    "mlp_fc1 / fc2  — İleri beslemeli bloktaki iki doğrusal katman.",
    "lm_head  — Son gizli durumu vocab boyutuna yansıtır → logitler.",
    "Tüm ağırlıklar küçük rastgele Gauss değerleriyle başlar (std=0.08). Eğitim onları şekillendirir.",
], Inches(0.6), Inches(1.1), Inches(7.0), Inches(3.2), size=15)

add_image(sl, img('initial_weights.png'), Inches(7.5), Inches(1.0), height=Inches(5.6))
label(sl, "Başlangıç ağırlıkları — saf gürültü.\nHer hücre bir ağırlık değeri.\nKırmızı = negatif, Mavi = pozitif.",
      Inches(7.5), Inches(6.6), Inches(5.5), Inches(0.7), size=11, color=MUTED, italic=True)

# ── 7. MİMARİ GENEL BAKIŞ ───────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Adım 5 — Model Mimarisi (Transformer)",
      Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=34, bold=True, color=GOLD)
divider(sl, Inches(0.95))

blocks = [
    ("Token + Pozisyon Gömme",
     "Harfin vektörünü (wte) ve konumun vektörünü (wpe) topla.\n"
     "Sonuç: hem NEYIN hem NEREDE olduğunu kodlayan 16 boyutlu bir vektör."),
    ("RMSNorm",
     "Vektörün karekök-ortalama-karesini 1'e normalize et.\n"
     "Sayıları ağ boyunca dengeli tutar."),
    ("Çok Başlı Öz-Dikkat",
     "Her konum önceki tüm konumlara bakar ve ne alacağına karar verir.\n"
     "4 kafa bunu paralelde yapar; her kafa farklı kalıplara odaklanır."),
    ("Artık Bağlantı",
     "Dikkat çıktısını girdiyle topla (atlama bağlantısı).\n"
     "Gradyanların erken katmanlara doğrudan akmasını sağlar — eğitimi kararlı kılar."),
    ("MLP Bloğu",
     "ReLU aktivasyonlu iki doğrusal katman — genişlik 4× büyür, sonra küçülür.\n"
     "Dikkatın bulduğu kalıpları saklar ve dönüştürür."),
    ("LM Kafası → Logitler → Softmax",
     "vocab_size kadar sayıya (logit) yansıt, olasılıklara çevir.\n"
     "En yüksek olasılıklı token modelin en iyi tahminidir."),
]
for i, (baslik, aciklama) in enumerate(blocks):
    top = Inches(1.1) + i * Inches(1.03)
    box(sl, Inches(0.4), top, Inches(0.08), Inches(0.78), GOLD)
    label(sl, baslik, Inches(0.6), top, Inches(4.5), Inches(0.38),
          size=14, bold=True, color=GOLD)
    label(sl, aciklama, Inches(0.6), top + Inches(0.37), Inches(12.2), Inches(0.55),
          size=13, color=LIGHT)

# ── 8. DİKKAT MEKANİZMASI ───────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Dikkat Mekanizması — Q, K, V",
      Inches(0.6), Inches(0.2), Inches(9), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))

label(sl, "Modelin içinde bir arama motoru gibi düşün:",
      Inches(0.6), Inches(1.1), Inches(12), Inches(0.45), size=18, color=LIGHT, bold=True)

qkv = [
    ("Q  Sorgu (Query)", GOLD,
     "\"Ne arıyorum?\"\nMevcut token bir soru sorar."),
    ("K  Anahtar (Key)",  BLUE,
     "\"Bende ne var?\"\nGeçmişteki her token kendini tanıtır."),
    ("V  Değer (Value)", GREEN,
     "\"Sana ne vereceğim?\"\nDikkat yüksekse aktarılacak gerçek bilgi."),
]
for i, (ad, renk, aciklama) in enumerate(qkv):
    left = Inches(0.5) + i * Inches(4.25)
    box(sl, left, Inches(1.7), Inches(3.9), Inches(1.8), ACCENT)
    label(sl, ad, left + Inches(0.15), Inches(1.8), Inches(3.6), Inches(0.5),
          size=20, bold=True, color=renk)
    label(sl, aciklama, left + Inches(0.15), Inches(2.3), Inches(3.6), Inches(1.1),
          size=14, color=LIGHT)

label(sl, "Skor  =  (Q · K) / √kafa_boyutu        →       softmax       →       V'nin ağırlıklı toplamı",
      Inches(0.6), Inches(3.65), Inches(12), Inches(0.5),
      size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

bullet_box(sl, [
    "Q·K iç çarpımı bir sorgunun her anahtarla ne kadar uyuştuğunu ölçer — yüksek = daha fazla dikkat.",
    "√kafa_boyutuna bölmek softmax'tan önce değerlerin çok büyümesini önler.",
    "Softmax skorları toplamı 1 olan olasılıklara çevirir.",
    "Her konum yalnızca kendine ve önceki konumlara bakabilir (nedensel maskeleme) — geleceği göremez.",
    "4 kafayı paralelde çalıştırmak modelin aynı anda farklı şeylere dikkat etmesini sağlar.",
], Inches(0.6), Inches(4.25), Inches(12.1), Inches(2.8), size=15)

# ── 9. DİKKAT ISI HARİTASI ──────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Dikkat Mekanizması Görselleştirildi",
      Inches(0.6), Inches(0.2), Inches(7), Inches(0.7),
      size=34, bold=True, color=GOLD)
divider(sl, Inches(0.95))

label(sl, "Aşağıdaki ısı haritası, eğitim boyunca bir kelime için dikkat ağırlıklarını gösteriyor.\n"
          "Satır = hangi konum bakıyor (sorgu).   Sütun = hangi konuma bakılıyor (anahtar).\n"
          "Daha parlak = daha fazla dikkat.",
      Inches(0.6), Inches(1.1), Inches(6.2), Inches(1.4), size=15, color=LIGHT)
label(sl, "Nelere dikkat etmeli:",
      Inches(0.6), Inches(2.6), Inches(6.2), Inches(0.35), size=15, bold=True, color=GOLD)
bullet_box(sl, [
    "Adım 1: düzgün bulanıklık — model her yere eşit bakar (hiçbir şey bilmiyor).",
    "Eğitim ortası: sol-alt üçgen keskinleşir — nedensel yapı ortaya çıkar.",
    "Geç eğitim: belirli hücreler parlar — hangi harflerin hangilerini öngördüğü öğrenilir.",
    "Köşegen neredeyse her zaman parlaktır: her harf en çok kendine dikkat eder.",
], Inches(0.6), Inches(3.0), Inches(6.2), Inches(2.8), size=14)
add_image(sl, img('attention_animated.gif'), Inches(6.8), Inches(1.0), height=Inches(6.0))

# ── 10. EĞİTİM DÖNGÜSÜ ──────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Adım 6 — Eğitim Döngüsü", Inches(0.6), Inches(0.2), Inches(9), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))

asama = [
    ("① İleri Geçiş",
     "Tokenize edilmiş bir ismi modele ver.\nHer konumda bir sonraki harf için olasılık dağılımı al."),
    ("② Kayıp",
     "Modelin doğru bir sonraki harfe ne kadar şaşırdığını ölç.\n"
     "Kayıp = −log(doğru tokene verilen olasılık).  Düşük = daha iyi."),
    ("③ Geri Geçiş",
     "Kayıp üzerinde .backward() çağır.\n"
     "Autograd hesap grafiğini dolaşır ve her ağırlık için .grad değerini doldurur."),
    ("④ Adam Güncellemesi",
     "Her ağırlığı kaybı azaltan yönde küçük bir adım kaydır.\n"
     "Adam, momentum tutarak her ağırlık için adım boyutunu ayrı ayrı ayarlar."),
    ("⑤ Gradyanları Sıfırla",
     "Bir sonraki adımdan önce tüm ağırlıkların .grad değerini 0'a sıfırla.\n"
     "Sıfırlanmazsa gradyanlar adımlar arasında birikerek yanlış sonuçlar verir."),
]
for i, (baslik, aciklama) in enumerate(asama):
    top = Inches(1.1) + i * Inches(1.15)
    box(sl, Inches(0.4), top, Inches(2.7), Inches(0.9), ACCENT)
    label(sl, baslik, Inches(0.5), top + Inches(0.15), Inches(2.5), Inches(0.6),
          size=15, bold=True, color=GOLD)
    label(sl, aciklama, Inches(3.3), top + Inches(0.08), Inches(9.6), Inches(0.85),
          size=14, color=LIGHT)

# ── 11. ADAM OPTİMİZÖRÜ ─────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Adam Optimizörü", Inches(0.6), Inches(0.2), Inches(9), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))

label(sl, "Düz gradyan inişi: ağırlık -= lr × grad\n"
          "Sorun: bazı ağırlıklar büyük adım ister, bazıları küçük. Tek lr ikisini birden yapamaz.\n"
          "Adam bunu her ağırlık için ayrı ve uyarlanabilir bir adım boyutu tutarak çözer.",
      Inches(0.6), Inches(1.1), Inches(12), Inches(1.1), size=16, color=LIGHT)

code_box(sl, """
# Her t adımında, i. ağırlık için:
m[i] = beta1 * m[i] + (1 - beta1) * grad      # momentum  (yumuşatılmış gradyan)
v[i] = beta2 * v[i] + (1 - beta2) * grad**2   # RMS       (yumuşatılmış kare gradyan)

m_hat = m[i] / (1 - beta1**t)   # tarafsızlaştırılmış (eğitim başında önemli)
v_hat = v[i] / (1 - beta2**t)

agirlik -= lr * m_hat / (sqrt(v_hat) + eps)    # her ağırlığa özel adım boyutu
""", Inches(0.6), Inches(2.3), Inches(12.1), Inches(2.8), size=14)

bullet_box(sl, [
    "m  momentum tamponu — gürültülü gradyanları yumuşatır.",
    "v  RMS tamponu — büyük gradyan alan ağırlıklar otomatik olarak daha küçük adım atar.",
    "beta1=0.85 ve beta2=0.99 tamponların güncelleme hızını kontrol eder.",
    "eps=1e-8 sıfıra bölmeyi önler.",
    "lr doğrusal olarak sıfıra iner — model savrulmak yerine yerleşir.",
], Inches(0.6), Inches(5.2), Inches(12.1), Inches(2.0), size=15)

# ── 12. EĞİTİM EĞRİLERİ ─────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Modelin Öğrenişini İzle", Inches(0.6), Inches(0.2), Inches(9), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))

label(sl, "Üstte: 1.000 adım boyunca eğitim kaybı.   Altta: öğrenme hızı (doğrusal azalma).",
      Inches(0.6), Inches(1.1), Inches(12), Inches(0.4), size=15, color=LIGHT)
add_image(sl, img('training_curves_animated.gif'), Inches(0.4), Inches(1.55), height=Inches(4.5))

bullet_box(sl, [
    "Kayıp yüksek başlar (~3) — model rastgele tahmin yapıyor.",
    "İlk başta hızlı düşer, sonra yavaşlar — kolay kalıplar önce öğrenilir.",
    "Bu görev için iyi bir son kayıp ~1.5–2.0 civarındadır.",
    "Öğrenme hızı sıfıra iner — güncellemeler eğitim bitince küçülür.",
], Inches(7.5), Inches(2.0), Inches(5.5), Inches(3.5), size=15)

label(sl, "Kayıp = modelin ne kadar şaşırdığı.\n"
          "Düşük = daha güvenli = daha iyi tahmin.",
      Inches(7.5), Inches(5.5), Inches(5.5), Inches(0.9),
      size=14, color=MUTED, italic=True)

# ── 13. AĞIRLIK DAĞILIMI ────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Ağırlık Dağılımı — Gürültüden Anlama",
      Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=34, bold=True, color=GOLD)
divider(sl, Inches(0.95))
add_image(sl, img('weight_dist_animated.gif'), Inches(0.3), Inches(1.1), height=Inches(5.5))
bullet_box(sl, [
    "Adım 1: geniş çan eğrisi — ağırlıklar Gauss(0, 0.08)'den çekiliyor.",
    "Eğitim: dağılım daralır ve kayar — işe yarayan ağırlıklar büyür, işe yaramayanlar küçülür.",
    "Ortalama sıfıra yakın kalır; yayılım (σ) modelin ne kadar 'kararlı' olduğunu gösterir.",
    "Çok geniş = patlayan gradyanlar; sıfıra çökmüş = ölü ağırlıklar — ikisi de kötü.",
    "initial_weights.png ve final_weights.png'yi karşılaştır — yapı oluştuğunu görebilirsin.",
], Inches(7.5), Inches(1.5), Inches(5.6), Inches(4.0), size=15)

# ── 14. BAŞLANGIÇ / SON AĞIRLIKLAR ──────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Başlangıç ve Son Ağırlıklar", Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))
label(sl, "Eğitimden önce", Inches(0.8), Inches(1.1), Inches(5.8), Inches(0.4),
      size=18, bold=True, color=MUTED)
label(sl, "Eğitimden sonra", Inches(7.1), Inches(1.1), Inches(5.8), Inches(0.4),
      size=18, bold=True, color=GREEN)
add_image(sl, img('initial_weights.png'), Inches(0.3), Inches(1.5), width=Inches(6.3))
add_image(sl, img('final_weights.png'),   Inches(6.8), Inches(1.5), width=Inches(6.3))
label(sl, "Rastgele gürültü — hiçbir yapı yok.",
      Inches(0.5), Inches(6.4), Inches(6.0), Inches(0.5), size=13, color=MUTED, italic=True)
label(sl, "Yapısal kalıplar — her matris bir rol üstlendi.",
      Inches(7.0), Inches(6.4), Inches(6.0), Inches(0.5), size=13, color=GREEN, italic=True)

# ── 15. GÖMME UZAYI PCA ─────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Model Ne Öğrendi? — Gömme Uzayı",
      Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=32, bold=True, color=GOLD)
divider(sl, Inches(0.95))

label(sl, "Her harfin 16 boyutlu öğrenilmiş bir vektörü var (wte).\n"
          "PCA bu 16 boyutu 2'ye indirir, böylece çizebiliriz.\n"
          "Eğitim ilerledikçe kümelerin nasıl oluştuğunu izle.",
      Inches(0.6), Inches(1.1), Inches(6.0), Inches(1.2), size=15, color=LIGHT)
bullet_box(sl, [
    "Kırmızı noktalar = ünlüler (a, e, i, o, u)",
    "Mavi noktalar = ünsüzler",
    "Yeşil yıldız = <BOS> tokeni",
    "Adım 1: her yer rastgele dağılmış — model hiçbir şey bilmiyor.",
    "Eğitim ortası: ünlüler ve ünsüzler birbirinden ayrılmaya başlar.",
    "Geç eğitim: benzer sesler kümelenir — b/p, m/n, a/e vb.",
    "<BOS> tüm harflerden uzaklaşır — eşsiz bir rolü var (başlangıç + bitiş).",
    "Bu modelin öğrendiği iç 'anlam sözlüğünün' görünür hale gelmesidir.",
], Inches(0.6), Inches(2.35), Inches(5.8), Inches(4.4), size=14)
add_image(sl, img('embedding_pca_animated.gif'), Inches(6.6), Inches(1.0), height=Inches(6.1))

# ── 16. SON PARAMETRE DAĞILIMI ───────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Son Parametre Dağılımı",
      Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))
add_image(sl, img('final_parameter_distribution.png'), Inches(0.3), Inches(1.1), height=Inches(5.0))
bullet_box(sl, [
    "Sol: son ağırlık değerleri — başlangıç Gauss'undan daha dar.",
    "Sağ: son gradyan değerleri — sıfıra yakın demek model yakınsadı demektir.",
    "Eğitim sonunda büyük gradyanlar = model hâlâ çok değişiyor (yakınsamadı).",
    "Kırmızı kesikli çizgi ortalamayı gösterir — 0'a yakın = sistematik yanlılık yok.",
], Inches(7.8), Inches(2.0), Inches(5.3), Inches(3.5), size=15)

# ── 17. ÇIKARIM ─────────────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Adım 7 — Çıkarım: İsim Üretme",
      Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=34, bold=True, color=GOLD)
divider(sl, Inches(0.95))

adimlar = [
    ("Başla",    "İlk token olarak <BOS> ver.  Model henüz hiçbir şey görmedi."),
    ("Tahmin Et","Model 27 token üzerinde bir olasılık dağılımı çıkarır."),
    ("Örnekle",  "Bir sonraki tokenı bu olasılıklarla ağırlıklı rastgele seç."),
    ("Ekle",     "O token bir sonraki giriş olur.  Bağlam penceresi büyür."),
    ("Tekrarla", "Model tekrar <BOS> çıkarana dek (= ismin sonu) devam et."),
]
for i, (baslik, aciklama) in enumerate(adimlar):
    top = Inches(1.2) + i * Inches(1.05)
    box(sl, Inches(0.4), top, Inches(1.7), Inches(0.8), ACCENT)
    label(sl, baslik, Inches(0.45), top + Inches(0.15), Inches(1.6), Inches(0.5),
          size=16, bold=True, color=GOLD)
    label(sl, aciklama, Inches(2.3), top + Inches(0.1), Inches(10.5), Inches(0.7),
          size=16, color=LIGHT)

label(sl, "Sıcaklık (Temperature) rastgeleliği kontrol eder:",
      Inches(0.6), Inches(6.45), Inches(12), Inches(0.35),
      size=15, bold=True, color=GOLD)
label(sl, "Düşük (0.1) = model her seferinde en olası harfi seçer → tekrarlı ama güvenli.   "
          "Yüksek (2.0) = çılgın tahminler → yaratıcı ama saçma.   0.5 bu görev için idealdir.",
      Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
      size=14, color=LIGHT)

# ── 18. ÜRETİM ANİMASYONU ───────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Üretim — Harf Harf",
      Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))

label(sl, "Çubuk grafik modelin her aday harf için olasılığını gösterir.\n"
          "Altın çubuk = gerçekte örneklenen harf.",
      Inches(0.6), Inches(1.1), Inches(6.0), Inches(0.9), size=15, color=LIGHT)
bullet_box(sl, [
    "İlk adımlar: olasılıklar çok harfe dağılmış — model hâlâ belirsiz.",
    "Birkaç harf sonra: dağılım keskinleşir — ismin kalıbı netleşiyor.",
    "Örn. 'Em' yazıldıktan sonra model güçlü biçimde 'm', 'i', 'a' tercih eder.",
    "Model hiçbir ismi ezberlemedi — kalıpları sıfırdan öğrendi.",
], Inches(0.6), Inches(2.1), Inches(6.0), Inches(3.0), size=14)
add_image(sl, img('generation_animated.gif'), Inches(6.5), Inches(1.0), height=Inches(6.1))

# ── 19. TEMEL DENKLEMLER ────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
label(sl, "Temel Denklemler — Tek Slayta Sığdırılmış",
      Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=34, bold=True, color=GOLD)
divider(sl, Inches(0.95))

denklemler = [
    ("Gömme",            "x  =  wte[token_id]  +  wpe[konum_id]"),
    ("RMSNorm",          "x̂  =  x / √(ortalama(x²) + ε)"),
    ("Dikkat skoru",     "skor(q, k)  =  (q · k) / √kafa_boyutu"),
    ("Softmax",          "softmax(zᵢ)  =  exp(zᵢ) / Σ exp(zⱼ)"),
    ("Dikkat çıktısı",   "çıktı  =  Σ softmax(skorlar)ₜ · vₜ"),
    ("Çapraz entropi",   "kayıp  =  −log p(doğru token)"),
    ("Adam adımı",       "w  −=  lr · m̂ / (√v̂ + ε)"),
]
for i, (ad, denklem) in enumerate(denklemler):
    top = Inches(1.1) + i * Inches(0.83)
    label(sl, ad, Inches(0.6), top, Inches(2.8), Inches(0.5),
          size=15, bold=True, color=GOLD)
    box(sl, Inches(3.5), top, Inches(9.4), Inches(0.65), ACCENT)
    label(sl, denklem, Inches(3.65), top + Inches(0.06), Inches(9.1), Inches(0.55),
          size=16, color=GREEN, bold=True)

# ── 20. ÖZET ────────────────────────────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
box(sl, 0, Inches(2.8), W, Inches(2.0), ACCENT)
label(sl, "Özet", Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
      size=36, bold=True, color=GOLD)
divider(sl, Inches(0.95))

cikarımlar = [
    "Bir GPT; matris çarpımları, softmax ve bir döngüden ibarettir.",
    "Autograd eğitimi mümkün kılar — elle türev almaya gerek yok.",
    "Dikkat mekanizması her konumun geçmişteki tüm konumlardan bilgi toplamasını sağlar.",
    "Gömmeler modelin öğrendiği anlam alfabesidir — PCA ile görünür hale gelir.",
    "Adam her ağırlık için adım boyutunu ayrı ayarlar — düz gradyan inişinden çok daha iyi.",
    "Sıcaklık tek başına 'kesin' ile 'yaratıcı' üretim arasındaki tek ayar noktasıdır.",
    "1.000 adım ve ~7.000 parametre, inandırıcı isimler hayal etmeye yeterlidir.",
]
for i, metin in enumerate(cikarımlar):
    top = Inches(1.1) + i * Inches(0.82)
    label(sl, f"{'★' if i < 3 else '·'}  {metin}",
          Inches(0.7), top, Inches(12), Inches(0.6),
          size=16, color=WHITE if i < 3 else LIGHT,
          bold=(i < 3))

label(sl, "microgpt.py  ·  ~380 satır  ·  sıfır bağımlılık  ·  eksiksiz bir GPT",
      Inches(0.6), Inches(6.8), Inches(12), Inches(0.45),
      size=13, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# ---------------------------------------------------------------------------
# Kaydet
# ---------------------------------------------------------------------------

cikis = os.path.join(HERE, 'microgpt_aciklamali.pptx')
prs.save(cikis)
print(f"Kaydedildi  {cikis}")
print(f"Slayt sayisi: {len(prs.slides)}")
